"""Build the Optimus Mexico overview PowerPoint deck.

Creates docs/Optimus_Mexico_Presentation.pptx: purpose, system flow,
one slide per module with explanation, app screenshots (when present in
docs/presentation_assets/), sample results and the technology stack.

Usage: python scripts/make_presentation.py
"""
from __future__ import annotations

from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
ASSETS = ROOT / "docs" / "presentation_assets"
OUT = ROOT / "docs" / "Optimus_Mexico_Presentation.pptx"

ACCENT = RGBColor(0x1F, 0x3A, 0x5F)
ACCENT_LIGHT = RGBColor(0x4A, 0x6F, 0xA5)
GREY = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _title_bar(slide, text: str):
    bar = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), SLIDE_W - Inches(1.0), Inches(0.9))
    tf = bar.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = ACCENT


def _bullets(slide, items: list[str], left, top, width, height, size=15):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if item.startswith("  "):
            p.level = 1
            p.text = item.strip()
            p.font.size = Pt(size - 2)
            p.font.color.rgb = GREY
        else:
            p.text = item
            p.font.size = Pt(size)
        p.space_after = Pt(6)
    return box


def _picture(slide, name: str, left, top, width=None, height=None) -> bool:
    path = ASSETS / name
    if not path.exists():
        ph = slide.shapes.add_textbox(left, top, width or Inches(5), height or Inches(3))
        ph.text_frame.text = f"[screenshot: {name}]"
        ph.text_frame.paragraphs[0].font.color.rgb = GREY
        ph.text_frame.paragraphs[0].font.size = Pt(12)
        return False
    slide.shapes.add_picture(str(path), left, top, width=width, height=height)
    return True


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def main() -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # ------------------------------------------------- 1. cover
    s = _blank(prs)
    bg = s.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)  # 1 = rectangle
    bg.fill.solid()
    bg.fill.fore_color.rgb = ACCENT
    bg.line.fill.background()
    box = s.shapes.add_textbox(Inches(1.0), Inches(2.4), SLIDE_W - Inches(2.0), Inches(2.5))
    tf = box.text_frame
    tf.text = "Optimus Mexico"
    tf.paragraphs[0].font.size = Pt(54)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    p = tf.add_paragraph()
    p.text = "PV + BESS revenue optimization for the Mexican electricity market"
    p.font.size = Pt(22)
    p.font.color.rgb = WHITE
    p = tf.add_paragraph()
    p.text = f"System overview & flow  |  {date.today():%B %Y}  |  Goldbeck Solar"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0xC9, 0xD6, 0xE8)

    # ------------------------------------------------- 2. purpose
    s = _blank(prs)
    _title_bar(s, "What is Optimus Mexico?")
    _bullets(s, [
        "Purpose: estimate how much money a solar + battery project would make at any node of the "
        "Mexican electricity market — and whether a behind-the-meter battery pays off for an "
        "industrial CFE customer.",
        "",
        "Two analysis tracks in one tool:",
        "  Wholesale (MEM): CENACE hourly PML prices + PV profile → optimized BESS dispatch → "
        "revenues from energy arbitrage, PPAs and the capacity market.",
        "  Behind-the-Meter (CFE): 15-minute load data → CFE bill reconstruction (GDMTH) → "
        "peak-shaving / TOU dispatch → investor decision with GO / REVISE / NO-GO.",
        "",
        "Deliverables: interactive web app, Excel / Word / PDF / PowerPoint reports, "
        "investor-grade summaries with CAPEX/OPEX/IRR, full pytest-verified engine (71 tests).",
        "Built to be integrated into Optimus AI after testing.",
    ], Inches(0.7), Inches(1.4), Inches(12.0), Inches(5.5), size=17)

    # ------------------------------------------------- 3. system flow (wholesale)
    s = _blank(prs)
    _title_bar(s, "System flow — wholesale market track")
    steps = [
        ("CENACE node catalog", "~3,000 NodosP enriched with zona, tipo, gerencia, estado"),
        ("Node Resolver", "pick nodes from dropdown or extract from PML reports"),
        ("Market & PV data", "8,760 h PML + PV profile (upload or MW AC/MWp/yield specs)"),
        ("BESS dispatch", "price-rank heuristic or daily linear optimization"),
        ("PPA / Capacity", "contract structures + 100 critical hours credit"),
        ("Reports", "Excel, Word, PDF, PowerPoint + investor report + project IRR"),
    ]
    top = Inches(1.5)
    box_h, gap = Inches(0.78), Inches(0.17)
    for i, (head, sub) in enumerate(steps):
        shp = s.shapes.add_shape(1, Inches(0.8), top + i * (box_h + gap), Inches(5.6), box_h)
        shp.fill.solid()
        shp.fill.fore_color.rgb = ACCENT_LIGHT if i % 2 == 0 else ACCENT
        shp.line.fill.background()
        tf = shp.text_frame
        tf.word_wrap = True
        tf.margin_top = Emu(40000)
        tf.text = head
        tf.paragraphs[0].font.size = Pt(15)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = WHITE
        p = tf.add_paragraph()
        p.text = sub
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(0xDD, 0xE6, 0xF2)
    _picture(s, "shot_02_flow.png", Inches(6.9), Inches(1.5), width=Inches(6.0))

    # ------------------------------------------------- 4. app home
    s = _blank(prs)
    _title_bar(s, "The app — horizontal navigation, four pages")
    _bullets(s, [
        "How it works — purpose, pipeline diagram, data inventory with live found/missing status.",
        "0. Node Resolver — searchable catalog dropdown with filters.",
        "1. Front-of-Meter Dispatch + PPA + Capacity — wholesale revenue + CAPEX/OPEX/IRR.",
        "2. Behind-the-Meter (CFE) — bill savings, readiness gate, savings IRR.",
        "PV input: upload a profile OR enter MW AC, MWp, yield and degradation.",
        "Launch: double-click run_app.bat (creates venv and installs dependencies automatically).",
    ], Inches(0.7), Inches(1.3), Inches(5.3), Inches(5.0), size=15)
    _picture(s, "shot_01_home.png", Inches(6.2), Inches(1.3), width=Inches(6.7))

    # ------------------------------------------------- 5. node resolver
    s = _blank(prs)
    _title_bar(s, "Module 0 — Node Resolver")
    _bullets(s, [
        "Every project starts at a node: the price point where energy is settled.",
        "Enriched CENACE catalog: 2,999 nodes with zona de carga, tipo "
        "(generación/carga), gerencia regional, estado and municipio.",
        "Pick nodes from a searchable dropdown (filter by zona, estado, tipo) — "
        "or extract node keys from a CENACE PML report (PDF/CSV/XLSX).",
        "Output: resolution workbook with full node context, ready for the "
        "dispatch module or the PML scraper.",
    ], Inches(0.7), Inches(1.3), Inches(5.3), Inches(5.0), size=15)
    _picture(s, "shot_03_resolver.png", Inches(6.2), Inches(1.3), width=Inches(6.7))

    # ------------------------------------------------- 6. dispatch
    s = _blank(prs)
    _title_bar(s, "Module 1 — Front-of-Meter Dispatch, PPA & Capacity")
    _bullets(s, [
        "Merges 8,760 h of nodal PML prices with PV (upload or synthesized from specs).",
        "Battery dispatch engines:",
        "  price_rank — transparent daily heuristic (charge cheap, discharge expensive)",
        "  lp — daily linear program with SOC carry-over, efficiency and cycle limits",
        "Duration sweep 2-8 h finds the most profitable battery size.",
        "PPA structures: pro-rata, baseload block, solar-only — vs merchant.",
        "Capacity market: 100 critical hours → accredited MW → MXN/year.",
        "Project economics: PV+BESS CAPEX, OPEX/insurance, financing → unlevered & levered IRR.",
    ], Inches(0.7), Inches(1.3), Inches(5.3), Inches(5.2), size=14)
    _picture(s, "shot_04_dispatch.png", Inches(6.2), Inches(1.3), width=Inches(6.7))

    # ------------------------------------------------- 7. revenue
    s = _blank(prs)
    _title_bar(s, "Module 1 — Revenue & project economics")
    _bullets(s, [
        "Interactive charts: average-day dispatch profile, single-day detail "
        "(flows + SOC + price), annual revenue by stream, monthly stacked revenue.",
        "Project economics panel: CAPEX (PV USD/kWp + BESS USD/kWh), OPEX, insurance, "
        "financing, project life → NPV, payback, unlevered and levered project IRR.",
        "Sample-data result: MXN 238M total revenue (PPA + Capacity + Merchant).",
        "One click exports: Excel workbook, Word, PDF, PowerPoint — plus investor report PDF.",
    ], Inches(0.7), Inches(1.3), Inches(5.3), Inches(5.0), size=15)
    _picture(s, "shot_05_revenue.png", Inches(6.2), Inches(1.3), width=Inches(6.7))

    # ------------------------------------------------- 8. BTM flow
    s = _blank(prs)
    _title_bar(s, "Module 2 — Behind-the-Meter (CFE) — how it flows")
    steps = [
        ("15-minute load + PV data", "upload or PV specs (MW AC / MWp / yield)"),
        ("CFE tariff calendar", "GDMTH base / intermedia / punta, seasons, holidays"),
        ("Baseline bill reconstruction", "energy by period + capacity & distribution demand"),
        ("Dispatch engine", "rule-based screening OR LP bankable (import-cap sweep)"),
        ("Optimized bill + savings", "savings by CFE tariff component; BTM-only revenues"),
        ("Investor layer", "readiness gate DEMO→INVESTMENT READY; IRR on CFE savings"),
    ]
    top = Inches(1.5)
    for i, (head, sub) in enumerate(steps):
        shp = s.shapes.add_shape(1, Inches(0.8), top + i * (box_h + gap), Inches(5.8), box_h)
        shp.fill.solid()
        shp.fill.fore_color.rgb = ACCENT if i % 2 == 0 else ACCENT_LIGHT
        shp.line.fill.background()
        tf = shp.text_frame
        tf.word_wrap = True
        tf.margin_top = Emu(40000)
        tf.text = head
        tf.paragraphs[0].font.size = Pt(15)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = WHITE
        p = tf.add_paragraph()
        p.text = sub
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(0xDD, 0xE6, 0xF2)
    _picture(s, "shot_06_btm_results.png", Inches(7.1), Inches(1.5), width=Inches(5.8))

    # ------------------------------------------------- 9. BTM dispatch detail
    s = _blank(prs)
    _title_bar(s, "Module 2 — Daily dispatch & bill savings")
    _bullets(s, [
        "The battery keeps grid import under a monthly cap (found by binary search) "
        "and shifts punta-period energy to cheap base hours.",
        "PV serves load first; surplus charges the battery; nothing is exported.",
        "Sample plant (956 kW peak, 500 kWp PV, 300 kW / 600 kWh BESS):",
        "  CFE bill before: MXN 11.67M/yr → after: MXN 8.68M/yr",
        "  Annual savings: MXN 2.99M | peak shaved 956 → 845 kW | zero export",
    ], Inches(0.7), Inches(1.3), Inches(5.3), Inches(5.0), size=15)
    _picture(s, "shot_07_btm_day.png", Inches(6.2), Inches(1.3), width=Inches(6.7))

    # ------------------------------------------------- 10. investor layer
    s = _blank(prs)
    _title_bar(s, "Module 2 — Investor layer & savings IRR")
    _bullets(s, [
        "Converts the technical model into an investment answer for an owner-manager.",
        "Explicit costs: CAPEX (PV USD/kWp + BESS USD/kWh × FX), OPEX %/yr, insurance %/yr, "
        "financing (cash / lease / loan) and configurable BESS life (default 20 years).",
        "IRR on CFE savings — unlevered (bill savings − OPEX vs CAPEX) and levered "
        "(bankable savings after financing); NPV and payback over BESS life.",
        "Investment readiness gate: DEMO | SCREENING | REVISE | INVESTMENT READY — "
        "GO/REVISE/NO-GO blocked unless INVESTMENT READY (sample data = DEMO).",
        "Conservative haircuts on bankable savings; red flags; Excel investor package.",
        "Sub-workflows: Screening, Bankable Investor, Customer Proposal, Technical Dispatch.",
    ], Inches(0.7), Inches(1.3), Inches(5.3), Inches(5.4), size=14)
    _picture(s, "shot_08_btm_investor.png", Inches(6.2), Inches(1.3), width=Inches(6.7))

    # ------------------------------------------------- 11. tech + status
    s = _blank(prs)
    _title_bar(s, "Technology, quality & next steps")
    _bullets(s, [
        "Stack: Python, Streamlit, pandas, scipy (HiGHS LP), Altair charts, "
        "openpyxl / python-docx / reportlab / python-pptx for reports.",
        "Agent architecture: each capability is an isolated, testable module — "
        "ready to plug into Optimus AI.",
        "Quality: 71 automated tests (wholesale dispatch, BTM tariff/bill/dispatch, "
        "BTM revenue guard, readiness gate, LP optimizer, PV synthesis, FOM economics).",
        "New agents: pv_synthesis, fom_investor, btm_lp_optimizer, btm_revenue_guard, "
        "btm_investment_readiness.",
        "Data: official CENACE node catalog (auto-download + enrich), CNE 2026 "
        "starter tariffs, synthetic sample profiles for instant demos.",
        "Next: CENACE PML live scraping hardening, CFE/CNE tariff scrapers, PDF bill "
        "parser, DIST/DIT tariffs, multi-year PPA cash flows, Optimus AI integration.",
    ], Inches(0.7), Inches(1.4), Inches(12.0), Inches(5.2), size=16)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    have = sum(1 for f in ["shot_01_home.png", "shot_02_flow.png", "shot_03_resolver.png",
                           "shot_04_dispatch.png", "shot_05_revenue.png", "shot_06_btm_results.png",
                           "shot_07_btm_day.png", "shot_08_btm_investor.png"] if (ASSETS / f).exists())
    print(f"Wrote {OUT} ({len(prs.slides.__iter__.__self__._sldIdLst)} slides, {have}/8 screenshots embedded)")


if __name__ == "__main__":
    main()
