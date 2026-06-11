"""Printout agent: Excel, Word, PDF and PowerPoint deliverables.

All document writers consume the same `ReportData` bundle so every format
tells the same story: project summary, scenario comparison, monthly
results, PPA split and capacity credit.
"""
from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path

import pandas as pd


@dataclass
class ReportData:
    project_name: str
    dispatch: pd.DataFrame
    scenario_summary: pd.DataFrame | None = None
    ppa_summary: pd.DataFrame | None = None
    capacity_summary: pd.DataFrame | None = None
    assumptions: dict = field(default_factory=dict)

    def monthly(self) -> pd.DataFrame:
        monthly = self.dispatch.copy()
        monthly["month"] = pd.to_datetime(monthly["datetime"]).dt.month
        agg = {
            "pml": "mean",
            "pv_mwh": "sum",
            "pv_to_grid_mwh": "sum",
            "pv_to_bess_mwh": "sum",
            "bess_discharge_mwh": "sum",
            "merchant_revenue": "sum",
        }
        agg = {k: v for k, v in agg.items() if k in monthly.columns}
        return monthly.groupby("month", as_index=False).agg(agg)

    def headline_metrics(self) -> list[tuple[str, str]]:
        d = self.dispatch
        delivered = d["pv_to_grid_mwh"].sum() + d["bess_discharge_mwh"].sum()
        metrics = [
            ("Average PML (MXN/MWh)", f"{d['pml'].mean():,.2f}"),
            ("PV generation (MWh)", f"{d['pv_mwh'].sum():,.0f}"),
            ("Energy delivered (MWh)", f"{delivered:,.0f}"),
            ("BESS discharge (MWh)", f"{d['bess_discharge_mwh'].sum():,.0f}"),
            ("Merchant revenue (MXN)", f"{d['merchant_revenue'].sum():,.0f}"),
        ]
        if self.capacity_summary is not None:
            cap = dict(zip(self.capacity_summary["metric"], self.capacity_summary["value"]))
            if "Capacity revenue (MXN/year)" in cap:
                metrics.append(("Capacity revenue (MXN/year)", f"{float(cap['Capacity revenue (MXN/year)']):,.0f}"))
        return metrics

    def revenue_streams(self) -> dict[str, float]:
        """Annual revenue by stream (MXN): PPA, merchant and capacity."""
        d = self.dispatch
        streams: dict[str, float] = {}
        if "ppa_revenue" in d.columns:
            streams["PPA"] = float(d["ppa_revenue"].sum())
        streams["Merchant"] = float(d["merchant_revenue"].sum())
        if self.capacity_summary is not None:
            cap = dict(zip(self.capacity_summary["metric"], self.capacity_summary["value"]))
            if "Capacity revenue (MXN/year)" in cap:
                streams["Capacity"] = float(cap["Capacity revenue (MXN/year)"])
        return streams

    def average_day(self) -> pd.DataFrame:
        """Mean hourly dispatch profile across the whole series."""
        d = self.dispatch.copy()
        d["hour"] = pd.to_datetime(d["datetime"]).dt.hour
        agg = {
            "pv_to_grid_mwh": "mean",
            "pv_to_bess_mwh": "mean",
            "bess_discharge_mwh": "mean",
            "pml": "mean",
        }
        agg = {k: v for k, v in agg.items() if k in d.columns}
        return d.groupby("hour", as_index=False).agg(agg)


# ---------------------------------------------------------------- Excel

def write_dispatch_workbook(dispatch: pd.DataFrame, scenario_summary: pd.DataFrame | None, out_path: str | Path) -> Path:
    out_path = Path(out_path)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    with pd.ExcelWriter(out_path, engine="openpyxl") as writer:
        dispatch.to_excel(writer, sheet_name="Dispatch_8760", index=False)
        monthly = dispatch.copy()
        monthly["month"] = pd.to_datetime(monthly["datetime"]).dt.month
        agg = {
            "pml": "mean",
            "pv_mwh": "sum",
            "pv_to_grid_mwh": "sum",
            "pv_to_bess_mwh": "sum",
            "bess_discharge_mwh": "sum",
            "merchant_revenue": "sum",
        }
        agg = {k: v for k, v in agg.items() if k in monthly.columns}
        monthly.groupby("month", as_index=False).agg(agg).to_excel(writer, sheet_name="Monthly_Summary", index=False)
        if scenario_summary is not None:
            scenario_summary.to_excel(writer, sheet_name="Scenario_Comparison", index=False)
    return out_path


def write_full_workbook(report: ReportData, out_path: str | Path) -> Path:
    out_path = Path(out_path)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    with pd.ExcelWriter(out_path, engine="openpyxl") as writer:
        pd.DataFrame(report.headline_metrics(), columns=["metric", "value"]).to_excel(writer, sheet_name="Summary", index=False)
        if report.assumptions:
            pd.DataFrame(list(report.assumptions.items()), columns=["assumption", "value"]).to_excel(writer, sheet_name="Assumptions", index=False)
        report.dispatch.to_excel(writer, sheet_name="Dispatch_8760", index=False)
        report.monthly().to_excel(writer, sheet_name="Monthly_Summary", index=False)
        if report.scenario_summary is not None:
            report.scenario_summary.to_excel(writer, sheet_name="Scenario_Comparison", index=False)
        if report.ppa_summary is not None:
            report.ppa_summary.to_excel(writer, sheet_name="PPA_Summary", index=False)
        if report.capacity_summary is not None:
            report.capacity_summary.to_excel(writer, sheet_name="Capacity_Summary", index=False)
    return out_path


# ---------------------------------------------------------------- chart

def _monthly_chart_png(report: ReportData, out_dir: Path) -> Path | None:
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
    except ImportError:
        return None
    monthly = report.monthly()
    if monthly.empty or "merchant_revenue" not in monthly.columns:
        return None
    fig, ax1 = plt.subplots(figsize=(9, 4.5))
    ax1.bar(monthly["month"], monthly["merchant_revenue"] / 1e6, color="#1f77b4", label="Merchant revenue")
    ax1.set_xlabel("Month")
    ax1.set_ylabel("Revenue (MXN millions)")
    if "pml" in monthly.columns:
        ax2 = ax1.twinx()
        ax2.plot(monthly["month"], monthly["pml"], color="#d62728", marker="o", label="Avg PML")
        ax2.set_ylabel("Avg PML (MXN/MWh)")
    fig.suptitle(f"{report.project_name} — Monthly results")
    fig.tight_layout()
    png = out_dir / "monthly_results.png"
    fig.savefig(png, dpi=150)
    plt.close(fig)
    return png


def _avg_day_chart_png(report: ReportData, out_dir: Path) -> Path | None:
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
    except ImportError:
        return None
    avg = report.average_day()
    if avg.empty:
        return None
    fig, ax1 = plt.subplots(figsize=(9, 4.5))
    bottom = None
    for col, label, color in [
        ("pv_to_grid_mwh", "PV to grid", "#f5b942"),
        ("pv_to_bess_mwh", "PV to BESS", "#7fb3d5"),
        ("bess_discharge_mwh", "BESS discharge", "#27ae60"),
    ]:
        if col not in avg.columns:
            continue
        vals = avg[col].to_numpy()
        ax1.bar(avg["hour"], vals, bottom=bottom, label=label, color=color)
        bottom = vals if bottom is None else bottom + vals
    ax1.set_xlabel("Hour of day")
    ax1.set_ylabel("Average MWh")
    ax1.legend(loc="upper left", fontsize=8)
    if "pml" in avg.columns:
        ax2 = ax1.twinx()
        ax2.plot(avg["hour"], avg["pml"], color="#d62728", marker="o", linewidth=1.5, label="Avg PML")
        ax2.set_ylabel("Avg PML (MXN/MWh)")
        ax2.legend(loc="upper right", fontsize=8)
    fig.suptitle(f"{report.project_name} — Average-day dispatch profile")
    fig.tight_layout()
    png = out_dir / "avg_day_dispatch.png"
    fig.savefig(png, dpi=150)
    plt.close(fig)
    return png


def _revenue_breakdown_png(report: ReportData, out_dir: Path) -> Path | None:
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
    except ImportError:
        return None
    streams = report.revenue_streams()
    if not streams:
        return None
    labels = list(streams.keys())
    values = [v / 1e6 for v in streams.values()]
    colors = {"PPA": "#1f77b4", "Merchant": "#f5b942", "Capacity": "#27ae60"}
    fig, ax = plt.subplots(figsize=(7, 4))
    bars = ax.bar(labels, values, color=[colors.get(l, "#888888") for l in labels])
    for bar, v in zip(bars, values):
        ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height(), f"{v:,.1f}", ha="center", va="bottom", fontsize=9)
    ax.set_ylabel("Annual revenue (MXN millions)")
    fig.suptitle(f"{report.project_name} — Revenue by stream")
    fig.tight_layout()
    png = out_dir / "revenue_breakdown.png"
    fig.savefig(png, dpi=150)
    plt.close(fig)
    return png


def _df_rows(df: pd.DataFrame, max_rows: int = 20) -> list[list[str]]:
    head = [str(c) for c in df.columns]
    body = [[(f"{v:,.2f}" if isinstance(v, float) else str(v)) for v in row] for row in df.head(max_rows).itertuples(index=False)]
    return [head] + body


# ---------------------------------------------------------------- Word

def write_word_report(report: ReportData, out_path: str | Path) -> Path:
    from docx import Document
    from docx.shared import Inches

    out_path = Path(out_path)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    doc = Document()
    doc.add_heading(f"{report.project_name} — PV+BESS Market Report", level=0)
    doc.add_paragraph("Automated output from Optimus Mexico. Results below are based on the configured PML series, PV 8760 profile and dispatch assumptions.")

    doc.add_heading("Headline metrics", level=1)
    table = doc.add_table(rows=0, cols=2)
    table.style = "Light Grid Accent 1"
    for name, value in report.headline_metrics():
        cells = table.add_row().cells
        cells[0].text = name
        cells[1].text = value

    if report.assumptions:
        doc.add_heading("Assumptions", level=1)
        table = doc.add_table(rows=0, cols=2)
        table.style = "Light Grid Accent 1"
        for k, v in report.assumptions.items():
            cells = table.add_row().cells
            cells[0].text = str(k)
            cells[1].text = str(v)

    def add_df(title: str, df: pd.DataFrame | None):
        if df is None or df.empty:
            return
        doc.add_heading(title, level=1)
        rows = _df_rows(df)
        table = doc.add_table(rows=0, cols=len(rows[0]))
        table.style = "Light Grid Accent 1"
        for row in rows:
            cells = table.add_row().cells
            for i, value in enumerate(row):
                cells[i].text = value

    add_df("Scenario comparison", report.scenario_summary)
    add_df("PPA vs merchant", report.ppa_summary)
    add_df("Capacity credit (100 critical hours)", report.capacity_summary)
    add_df("Monthly summary", report.monthly())

    chart = _monthly_chart_png(report, out_path.parent)
    if chart:
        doc.add_heading("Monthly results", level=1)
        doc.add_picture(str(chart), width=Inches(6.5))

    doc.save(out_path)
    return out_path


# ---------------------------------------------------------------- PDF

def write_pdf_report(report: ReportData, out_path: str | Path) -> Path:
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.lib.units import inch
    from reportlab.platypus import Image, Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle

    out_path = Path(out_path)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    styles = getSampleStyleSheet()
    story = [Paragraph(f"{report.project_name} — PV+BESS Market Report", styles["Title"]), Spacer(1, 12)]

    def add_table(title: str, rows: list[list[str]]):
        story.append(Paragraph(title, styles["Heading2"]))
        table = Table(rows, hAlign="LEFT")
        table.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#1f77b4")),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
            ("FONTSIZE", (0, 0), (-1, -1), 7),
            ("GRID", (0, 0), (-1, -1), 0.4, colors.grey),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#eef3f8")]),
        ]))
        story.append(table)
        story.append(Spacer(1, 12))

    add_table("Headline metrics", [["Metric", "Value"]] + [[n, v] for n, v in report.headline_metrics()])
    if report.assumptions:
        add_table("Assumptions", [["Assumption", "Value"]] + [[str(k), str(v)] for k, v in report.assumptions.items()])
    if report.scenario_summary is not None and not report.scenario_summary.empty:
        add_table("Scenario comparison", _df_rows(report.scenario_summary))
    if report.ppa_summary is not None and not report.ppa_summary.empty:
        add_table("PPA vs merchant", _df_rows(report.ppa_summary))
    if report.capacity_summary is not None and not report.capacity_summary.empty:
        add_table("Capacity credit (100 critical hours)", _df_rows(report.capacity_summary))
    add_table("Monthly summary", _df_rows(report.monthly()))

    chart = _monthly_chart_png(report, out_path.parent)
    if chart:
        story.append(Paragraph("Monthly results", styles["Heading2"]))
        story.append(Image(str(chart), width=6.5 * inch, height=3.25 * inch))

    SimpleDocTemplate(str(out_path), pagesize=letter).build(story)
    return out_path


# ---------------------------------------------------------------- PowerPoint

def write_ppt_report(report: ReportData, out_path: str | Path) -> Path:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    out_path = Path(out_path)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = f"{report.project_name}"
    slide.placeholders[1].text = "PV+BESS Market Report — Optimus Mexico"

    def add_table_slide(title: str, rows: list[list[str]]):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = title
        n_rows, n_cols = len(rows), len(rows[0])
        shape = slide.shapes.add_table(n_rows, n_cols, Inches(0.4), Inches(1.4), Inches(9.2), Inches(0.3 * n_rows))
        table = shape.table
        for r, row in enumerate(rows):
            for c, value in enumerate(row):
                cell = table.cell(r, c)
                cell.text = value
                for p in cell.text_frame.paragraphs:
                    for run in p.runs:
                        run.font.size = Pt(10)

    add_table_slide("Headline metrics", [["Metric", "Value"]] + [[n, v] for n, v in report.headline_metrics()])
    if report.scenario_summary is not None and not report.scenario_summary.empty:
        add_table_slide("Scenario comparison", _df_rows(report.scenario_summary, max_rows=8))
    if report.ppa_summary is not None and not report.ppa_summary.empty:
        add_table_slide("PPA vs merchant", _df_rows(report.ppa_summary, max_rows=8))
    if report.capacity_summary is not None and not report.capacity_summary.empty:
        add_table_slide("Capacity credit", _df_rows(report.capacity_summary, max_rows=8))

    chart = _monthly_chart_png(report, out_path.parent)
    if chart:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "Monthly results"
        slide.shapes.add_picture(str(chart), Inches(0.5), Inches(1.5), width=Inches(9.0))

    prs.save(out_path)
    return out_path


# ---------------------------------------------------------------- Investor report

def write_investor_report(report: ReportData, out_path: str | Path) -> Path:
    """Polished executive-summary PDF aimed at investors.

    One narrative document: executive summary, investment highlights,
    revenue-by-stream chart, average-day dispatch chart, monthly results,
    key assumptions and a disclaimer.
    """
    from datetime import date

    from reportlab.lib import colors
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.lib.units import inch
    from reportlab.platypus import Image, PageBreak, Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle

    out_path = Path(out_path)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    styles = getSampleStyleSheet()
    accent = colors.HexColor("#1f3a5f")
    h1 = ParagraphStyle("InvH1", parent=styles["Heading1"], textColor=accent)
    body = ParagraphStyle("InvBody", parent=styles["BodyText"], fontSize=10, leading=14)
    small = ParagraphStyle("InvSmall", parent=styles["BodyText"], fontSize=7.5, leading=10, textColor=colors.grey)

    d = report.dispatch
    delivered = float(d["pv_to_grid_mwh"].sum() + d["bess_discharge_mwh"].sum())
    streams = report.revenue_streams()
    total_revenue = sum(streams.values())
    capture = total_revenue / delivered if delivered else 0.0
    avg_pml = float(d["pml"].mean())

    story = [
        Paragraph(report.project_name, ParagraphStyle("Cover", parent=styles["Title"], fontSize=26, textColor=accent)),
        Paragraph("PV + BESS Investment Summary — Optimus Mexico", styles["Heading2"]),
        Paragraph(f"Prepared {date.today():%d %B %Y}", small),
        Spacer(1, 18),
        Paragraph("Executive summary", h1),
        Paragraph(
            f"This analysis models a solar + storage project settled at the Mexican wholesale "
            f"market (MEM). Over the modeled year the project delivers {delivered:,.0f} MWh "
            f"and generates total revenue of MXN {total_revenue:,.0f} "
            f"({total_revenue / 1e6:,.1f} million), equivalent to a capture price of "
            f"MXN {capture:,.0f}/MWh against an average nodal PML of MXN {avg_pml:,.0f}/MWh. "
            f"Revenue is built from {', '.join(f'{k.lower()} (MXN {v/1e6:,.1f}M)' for k, v in streams.items())}.",
            body,
        ),
        Spacer(1, 10),
        Paragraph("Investment highlights", h1),
    ]

    rows = [["Metric", "Value"]] + [[n, v] for n, v in report.headline_metrics()]
    rows.append(["Total revenue (MXN)", f"{total_revenue:,.0f}"])
    rows.append(["Capture price (MXN/MWh)", f"{capture:,.0f}"])
    table = Table(rows, hAlign="LEFT", colWidths=[3.2 * inch, 2.4 * inch])
    table.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), accent),
        ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
        ("FONTSIZE", (0, 0), (-1, -1), 9),
        ("GRID", (0, 0), (-1, -1), 0.4, colors.grey),
        ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#eef3f8")]),
        ("TOPPADDING", (0, 0), (-1, -1), 4),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
    ]))
    story += [table, Spacer(1, 14)]

    chart_rev = _revenue_breakdown_png(report, out_path.parent)
    if chart_rev:
        story += [Paragraph("Revenue by stream", h1), Image(str(chart_rev), width=5.6 * inch, height=3.2 * inch), Spacer(1, 10)]

    story.append(PageBreak())

    chart_day = _avg_day_chart_png(report, out_path.parent)
    if chart_day:
        story += [Paragraph("Average-day dispatch profile", h1), Image(str(chart_day), width=6.5 * inch, height=3.25 * inch), Spacer(1, 10)]

    chart_month = _monthly_chart_png(report, out_path.parent)
    if chart_month:
        story += [Paragraph("Monthly results", h1), Image(str(chart_month), width=6.5 * inch, height=3.25 * inch), Spacer(1, 10)]

    if report.assumptions:
        story.append(Paragraph("Key assumptions", h1))
        rows = [["Assumption", "Value"]] + [[str(k), str(v)] for k, v in report.assumptions.items()]
        table = Table(rows, hAlign="LEFT", colWidths=[3.2 * inch, 2.4 * inch])
        table.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), accent),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
            ("FONTSIZE", (0, 0), (-1, -1), 9),
            ("GRID", (0, 0), (-1, -1), 0.4, colors.grey),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#eef3f8")]),
        ]))
        story += [table, Spacer(1, 14)]

    story.append(Paragraph(
        "Disclaimer: this document is an automated modeling output based on historical or synthetic "
        "market data and the stated assumptions. It does not constitute investment advice, a revenue "
        "guarantee or an offer of securities. Actual results will differ.",
        small,
    ))

    SimpleDocTemplate(str(out_path), pagesize=letter, topMargin=0.7 * inch, bottomMargin=0.7 * inch).build(story)
    return out_path


def write_all_reports(report: ReportData, out_dir: str | Path, basename: str = "Optimus_Report", investor: bool = False) -> dict[str, Path]:
    """Write Excel + Word + PDF + PowerPoint (+ optional investor PDF). Skips formats whose library is missing."""
    out_dir = Path(out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    written: dict[str, Path] = {}
    written["xlsx"] = write_full_workbook(report, out_dir / f"{basename}.xlsx")
    for fmt, writer in [("docx", write_word_report), ("pdf", write_pdf_report), ("pptx", write_ppt_report)]:
        try:
            written[fmt] = writer(report, out_dir / f"{basename}.{fmt}")
        except ImportError as exc:
            print(f"Skipping {fmt}: {exc}")
    if investor:
        try:
            written["investor_pdf"] = write_investor_report(report, out_dir / f"{basename}_Investor.pdf")
        except ImportError as exc:
            print(f"Skipping investor report: {exc}")
    return written
